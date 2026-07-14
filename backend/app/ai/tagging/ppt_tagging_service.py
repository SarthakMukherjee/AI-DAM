from pptx import Presentation

class PPTTaggingService:
    def extract_text(self, ppt_path: str) -> str:
        try:
            prs = Presentation(ppt_path)
            full_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        full_text.append(shape.text.strip())
            return "\n".join(full_text).strip()
        except Exception as e:
            print(f"PPT text extraction failed: {e}")
            return ""

    def extract_metadata(self, ppt_path: str) -> dict:
        try:
            prs = Presentation(ppt_path)
            core_props = prs.core_properties
            return {
                "title": core_props.title or "",
                "author": core_props.author or "",
                "subject": core_props.subject or "",
                "keywords": core_props.keywords or ""
            }
        except Exception as e:
            print(f"PPT metadata extraction failed: {e}")
            return {
                "title": "",
                "author": "",
                "subject": "",
                "keywords": ""
            }
