# PPT PREVIEW SERVICE

import os
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont

from app.core.config.settings import settings


class PPTPreviewService:

    def generate_preview(
        self,
        ppt_path: str,
        filename: str
    ):
        """
        Extracts text from a PowerPoint presentation and renders a simple image preview.
        """
        # -----------------------------------
        # CREATE PREVIEW DIRECTORY
        # -----------------------------------
        preview_dir = os.path.join(
            settings.STORAGE_PATH,
            "previews",
            filename
        )
        os.makedirs(preview_dir, exist_ok=True)

        try:
            # -----------------------------------
            # EXTRACT TEXT
            # -----------------------------------
            prs = Presentation(ppt_path)
            full_text = []
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        full_text.append(shape.text.strip())
                if sum(len(t) for t in full_text) > 1500:
                    break
                    
            text_content = "\n".join(full_text)
            if not text_content.strip():
                text_content = "(Empty Presentation)"

            # -----------------------------------
            # RENDER TEXT TO IMAGE
            # -----------------------------------
            img_width, img_height = 800, 1000
            bg_color = (255, 255, 255)
            text_color = (0, 0, 0)
            
            image = Image.new("RGB", (img_width, img_height), bg_color)
            draw = ImageDraw.Draw(image)
            
            font = None
            try:
                font = ImageFont.truetype("arial.ttf", 16)
            except IOError:
                font = ImageFont.load_default()

            def wrap_text(text, font, max_width):
                lines = []
                for paragraph in text.split('\n'):
                    words = paragraph.split()
                    if not words:
                        lines.append("")
                        continue
                    current_line = []
                    for word in words:
                        current_line.append(word)
                        try:
                            w = draw.textlength(" ".join(current_line), font=font)
                        except AttributeError:
                            w, _ = draw.textsize(" ".join(current_line), font=font)
                            
                        if w > max_width:
                            current_line.pop()
                            if current_line:
                                lines.append(" ".join(current_line))
                            current_line = [word]
                    if current_line:
                        lines.append(" ".join(current_line))
                return lines

            wrapped_lines = wrap_text(text_content, font, img_width - 80)
            
            y_text = 40
            for line in wrapped_lines:
                if y_text > img_height - 40:
                    draw.text((40, y_text), "...", font=font, fill=text_color)
                    break
                draw.text((40, y_text), line, font=font, fill=text_color)
                
                try:
                    try:
                        bbox = draw.textbbox((0, 0), line, font=font)
                        line_height = bbox[3] - bbox[1]
                    except AttributeError:
                        _, line_height = draw.textsize(line, font=font)
                except Exception:
                    line_height = 20
                    
                y_text += line_height + 5

            page_filename = "page_1.png"
            page_path = os.path.join(preview_dir, page_filename)
            
            image.save(page_path)

            return page_path.replace("\\", "/")
            
        except Exception as e:
            print(f"[PPTPreviewService] Error generating preview: {e}")
            return None
